"""
File parser utilities
- Extract text from PDF, Word, PowerPoint
- Parse student lists from CSV/Excel
"""

import os
import re
from typing import List, Dict, Optional
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from io import BytesIO

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    
    # Method 1: pdfplumber (better for tables)
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error with pdfplumber: {e}")
    
    # Fallback: PyPDF2
    if not text.strip():
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error with PyPDF2: {e}")
    
    return clean_text(text)

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word document (.docx)"""
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return clean_text(text)
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint (.pptx)"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_text(text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def parse_student_list(file_path: str) -> List[Dict[str, str]]:
    """
    Parse student list from CSV or Excel file
    Expected columns: username, password, email (optional), full_name (optional)
    """
    students = []
    
    try:
        # Detect file type
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path)
        elif file_ext == '.csv':
            df = pd.read_csv(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        # Normalize column names
        df.columns = [col.strip().lower() for col in df.columns]
        
        # Required columns
        if 'username' not in df.columns:
            raise ValueError("File must contain 'username' column")
        
        # Optional columns with defaults
        for _, row in df.iterrows():
            student = {
                'username': str(row['username']).strip(),
                'password': str(row.get('password', '123456')).strip(),
                'email': str(row.get('email', f"{row['username']}@school.edu")).strip(),
                'full_name': str(row.get('full_name', row['username'])).strip(),
                'role': 'student'
            }
            students.append(student)
        
        return students
        
    except Exception as e:
        print(f"Error parsing student list: {e}")
        raise

def clean_text(text: str) -> str:
    """Clean and normalize extracted text"""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters but keep letters, numbers, and basic punctuation
    text = re.sub(r'[^\w\s.,!?;:\'"\-]', ' ', text)
    # Remove multiple spaces
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def detect_file_type(filename: str) -> str:
    """Detect file type from extension"""
    ext = os.path.splitext(filename)[1].lower()
    
    if ext in ['.pdf']:
        return 'pdf'
    elif ext in ['.docx', '.doc']:
        return 'docx'
    elif ext in ['.pptx', '.ppt']:
        return 'pptx'
    elif ext in ['.xlsx', '.xls']:
        return 'excel'
    elif ext == '.csv':
        return 'csv'
    else:
        return 'unknown'

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Extract text based on file type"""
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type in ['docx', 'doc']:
        return extract_text_from_docx(file_path)
    elif file_type in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

# Test function
if __name__ == "__main__":
    print("Testing file parsers...")
    
    # Test PDF
    test_pdf = "/tmp/test.pdf"
    if os.path.exists(test_pdf):
        text = extract_text_from_pdf(test_pdf)
        print(f"PDF extracted {len(text)} characters")
    
    # Test DOCX
    test_docx = "/tmp/test.docx"
    if os.path.exists(test_docx):
        text = extract_text_from_docx(test_docx)
        print(f"DOCX extracted {len(text)} characters")
    
    # Test PPTX
    test_pptx = "/tmp/test.pptx"
    if os.path.exists(test_pptx):
        text = extract_text_from_pptx(test_pptx)
        print(f"PPTX extracted {len(text)} characters")
    
    print("✓ File parsers ready")